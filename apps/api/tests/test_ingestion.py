"""Tests for context ingestion (issue #6)."""

from __future__ import annotations

from collections.abc import Iterator

import pytest
from fastapi.testclient import TestClient

from sanba_api.auth import create_session_token
from sanba_api.auth_google import AuthUser, require_user
from sanba_api.config import settings
from sanba_api.ingestion import (
    ContextIndexer,
    DocumentExtractionError,
    chunk_text,
    extract_text_from_upload,
)
from sanba_api.main import app

client = TestClient(app)


def _session_auth(session_id: str, role: str = "pm") -> dict[str, str]:
    """context 投稿は join 済みトークン必須（契約 §4）。テスト用の Bearer を作る。"""
    token = create_session_token(
        session_id, "owner-123456789", role, settings.session_signing_secret, 3600
    )
    return {"Authorization": f"Bearer {token}"}


@pytest.fixture(autouse=True)
def _assume_logged_in() -> Iterator[None]:
    """セッション作成に必要な検証済みユーザーをスタブする (ADR-0012)。"""
    app.dependency_overrides[require_user] = lambda: AuthUser(
        sub="owner-123456789", email="owner@example.com", email_verified=True, name="Owner"
    )
    yield
    app.dependency_overrides.pop(require_user, None)


def test_chunk_text_splits_on_paragraphs() -> None:
    text = "段落1の内容。\n\n段落2の内容。\n\n段落3の内容。"
    chunks = chunk_text(text, chunk_size=20)
    assert len(chunks) >= 2
    assert all(c.strip() for c in chunks)


def test_chunk_text_empty_returns_empty() -> None:
    assert chunk_text("   ") == []


def test_long_paragraph_is_windowed() -> None:
    chunks = chunk_text("あ" * 1000, chunk_size=300, overlap=50)
    assert len(chunks) > 1


def test_extract_text_from_txt_upload() -> None:
    assert extract_text_from_upload("notes.md", "# 見出し\n本文".encode()) == "# 見出し\n本文"


# ── 拡張形式のテキスト抽出（html/csv/json/docx/xlsx/pptx）───────────────────
def test_extract_text_from_html_strips_markup_and_scripts() -> None:
    html = (
        "<html><head><title>仕様</title><style>body{color:red}</style>"
        "<script>alert('x')</script></head>"
        "<body><h1>検索機能</h1><p>一覧を&amp;高速化する。</p></body></html>"
    )
    text = extract_text_from_upload("spec.html", html.encode())
    assert "検索機能" in text
    assert "一覧を&高速化する。" in text
    # マークアップ・script/style は本文に混ぜない。
    assert "alert" not in text
    assert "color:red" not in text
    assert "<h1>" not in text


def test_extract_text_from_csv_and_json_decodes_as_text() -> None:
    assert "画面,要件" in extract_text_from_upload("req.csv", "画面,要件\n検索,高速化".encode())
    assert '"goal"' in extract_text_from_upload("req.json", b'{"goal": "search"}')


def test_extract_text_from_docx_paragraphs_and_tables() -> None:
    import io

    from docx import Document

    document = Document()
    document.add_paragraph("要約機能が必要。")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "画面"
    table.rows[0].cells[1].text = "検索"
    buf = io.BytesIO()
    document.save(buf)

    text = extract_text_from_upload("spec.docx", buf.getvalue())
    assert "要約機能が必要。" in text
    assert "画面\t検索" in text


def test_extract_text_from_xlsx_all_sheets() -> None:
    import io

    from openpyxl import Workbook

    workbook = Workbook()
    first = workbook.active
    first.title = "要件"
    first.append(["画面", "要件"])
    first.append(["検索", "高速化"])
    second = workbook.create_sheet("課題")
    second.append(["優先度", "高"])
    buf = io.BytesIO()
    workbook.save(buf)

    text = extract_text_from_upload("req.xlsx", buf.getvalue())
    assert "# 要件" in text
    assert "検索\t高速化" in text
    # 先頭シートだけでなく全シートを読む。
    assert "# 課題" in text
    assert "優先度\t高" in text


def test_extract_text_from_pptx_slides_and_notes() -> None:
    import io

    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "検索リニューアル"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "対象は社内ユーザー"
    slide.notes_slide.notes_text_frame.text = "補足: 移行は段階的に"
    buf = io.BytesIO()
    presentation.save(buf)

    text = extract_text_from_upload("deck.pptx", buf.getvalue())
    assert "# スライド1" in text
    assert "検索リニューアル" in text
    assert "対象は社内ユーザー" in text
    assert "補足: 移行は段階的に" in text


def test_extract_text_selects_extractor_by_mime_when_no_extension() -> None:
    """拡張子なし・MIME のみのアップロードでも正しい抽出器を選ぶ（Codex P2）。

    受理判定（is_text_upload）は MIME だけでも通すため、抽出も MIME にフォールバック
    しないと ZIP バイト列を UTF-8 デコードして索引してしまう。
    """
    import io

    from docx import Document

    document = Document()
    document.add_paragraph("要約機能が必要。")
    buf = io.BytesIO()
    document.save(buf)

    text = extract_text_from_upload(
        "book",  # 拡張子なし
        buf.getvalue(),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    assert "要約機能が必要。" in text
    assert "PK" not in text  # ZIP バイト列をそのまま返していない。


def test_extract_text_from_broken_document_raises_typed_error() -> None:
    # 壊れたバイナリ文書は型付き例外で伝える（呼び出し側が 500 にせず「抽出 0 件」へ平しつつ、
    # メトリクスでは成功と区別して計上する）。
    for name in ("broken.docx", "broken.xlsx", "broken.pptx"):
        with pytest.raises(DocumentExtractionError):
            extract_text_from_upload(name, b"not-a-zip")


def test_zip_bomb_is_rejected_before_expansion(monkeypatch: pytest.MonkeyPatch) -> None:
    """展開後サイズが上限を超える zip コンテナは展開せずに弾く（OOM 防止）。"""
    import io

    from openpyxl import Workbook

    from sanba_api import ingestion

    workbook = Workbook()
    sheet = workbook.active
    for i in range(50):
        sheet.append([f"データ{i}" * 10])
    buf = io.BytesIO()
    workbook.save(buf)

    # 上限を人工的に絞り、正規の xlsx でも「展開後が大きすぎる」ケースを再現する。
    monkeypatch.setattr(ingestion, "_MAX_ZIP_EXPANSION_BYTES", 10)
    with pytest.raises(DocumentExtractionError):
        extract_text_from_upload("bomb.xlsx", buf.getvalue())


def test_upload_broken_docx_returns_zero_chunks() -> None:
    """壊れた文書のアップロードは 500 にせず indexed_chunks=0 で返る（best-effort）。"""
    created = client.post(
        "/api/sessions", json={"roles": ["pm"], "consent_acknowledged": True}
    ).json()
    sid = created["session_id"]
    res = client.post(
        f"/api/sessions/{sid}/context/file",
        files={"file": ("broken.docx", b"not-a-zip", "application/octet-stream")},
        headers=_session_auth(sid),
    )
    assert res.status_code == 200
    body = res.json()
    assert body["indexed_chunks"] == 0
    assert body["asset_kind"] == "doc"


def test_memory_indexer_counts_chunks() -> None:
    indexer = ContextIndexer()
    assert indexer.is_memory is True
    n = indexer.index_context("sess-1", ["a", "b", "c"], "spec.md")
    assert n == 3


# ── grounding 索引の取消（#245 真の破棄）────────────────────────────────────
def test_delete_context_removes_only_matching_source() -> None:
    """出所接頭辞 `asset:{id}` の chunk だけを消し、他 source/他セッションは残す。"""
    indexer = ContextIndexer()
    indexer.index_context("sess-1", ["o1", "o2"], "asset:asset-aaa")
    indexer.index_context("sess-1", ["k1"], "asset:asset-bbb")
    indexer.index_context("sess-1", ["t1"], "prd.md")
    indexer.index_context("sess-2", ["x1"], "asset:asset-aaa")  # 別セッションは無関係。

    removed = indexer.delete_context("sess-1", "asset:asset-aaa")
    assert removed == 2

    remaining = {(d["session_id"], d["source"]) for d in indexer._mem}
    # 対象セッションの asset-aaa#* は消える。
    assert ("sess-1", "asset:asset-aaa#0") not in remaining
    assert ("sess-1", "asset:asset-aaa#1") not in remaining
    # 別 asset・別 source・別セッションは残る（巻き込まない）。
    assert ("sess-1", "asset:asset-bbb#0") in remaining
    assert ("sess-1", "prd.md#0") in remaining
    assert ("sess-2", "asset:asset-aaa#0") in remaining


def test_delete_context_is_idempotent_when_absent() -> None:
    """存在しない出所の取消は 0 件で安全（冪等）。"""
    indexer = ContextIndexer()
    indexer.index_context("sess-1", ["o1"], "asset:asset-aaa")
    assert indexer.delete_context("sess-1", "asset:asset-zzz") == 0
    assert indexer.delete_context("sess-1", "asset:asset-aaa") == 1
    assert indexer.delete_context("sess-1", "asset:asset-aaa") == 0


def test_context_endpoint_indexes_text() -> None:
    created = client.post(
        "/api/sessions", json={"roles": ["pm"], "consent_acknowledged": True}
    ).json()
    sid = created["session_id"]
    res = client.post(
        f"/api/sessions/{sid}/context",
        json={"text": "要約機能が必要。\n\n対象は社内ユーザー。", "source_name": "prd.md"},
        headers=_session_auth(sid),
    )
    assert res.status_code == 200
    assert res.json()["indexed_chunks"] >= 1


def test_context_endpoint_requires_session_token() -> None:
    """匿名での RAG グラウンディング汚染を防ぐ（join 済みトークン必須）。"""
    res = client.post(
        "/api/sessions/sess-anon/context",
        json={"text": "注入テキスト", "source_name": "x"},
    )
    assert res.status_code == 401


def test_context_endpoint_rejects_token_for_other_session() -> None:
    """別セッションのトークンでは投稿できない（session_id 不一致）。"""
    res = client.post(
        "/api/sessions/sess-target/context",
        json={"text": "注入テキスト", "source_name": "x"},
        headers=_session_auth("sess-OTHER"),
    )
    assert res.status_code == 403


def test_context_endpoint_rejects_oversized() -> None:
    created = client.post(
        "/api/sessions", json={"roles": ["pm"], "consent_acknowledged": True}
    ).json()
    sid = created["session_id"]
    res = client.post(
        f"/api/sessions/{sid}/context",
        json={"text": "x" * 200_001, "source_name": "big.txt"},
        headers=_session_auth(sid),
    )
    assert res.status_code == 413
