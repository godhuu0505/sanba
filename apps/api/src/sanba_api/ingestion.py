"""Pre-interview context ingestion -> Elasticsearch grounding (issue #6).

Lets an owner register existing material (PRD drafts, meeting notes, specs) before
the interview. Chunks are indexed into the SAME Elasticsearch index the agent reads
via `search_grounding` (kind="context"), so the agent grounds its questions on the
user's real documents and skips already-answered topics.

NOTE: the document shape mirrors apps/agent/.../retrieval.py. A shared package is the
right long-term home (tracked as a follow-up); kept compact here on purpose.
"""

from __future__ import annotations

import io
import re
import zipfile
from collections.abc import Callable
from html.parser import HTMLParser

import structlog

from .config import settings
from .pii import mask_pii

log = structlog.get_logger(__name__)

INDEX = "sanba-grounding"
EMBED_DIM = 3072  # gemini-embedding-001 (must match apps/agent retrieval.py: same index)


def _source_matches(source: str, prefix: str) -> bool:
    """出所が削除接頭辞に一致するか（`prefix` 自身、または `prefix#<i>` の chunk）。

    `asset:{id}` を渡したとき `asset:{id}#0` 等を消しつつ、別 source（`asset:{id2}#0`）を
    巻き込まない（index_context が付ける `#<i>` 境界で判定する）。
    """
    return source == prefix or source.startswith(f"{prefix}#")


def chunk_text(text: str, chunk_size: int = 600, overlap: int = 80) -> list[str]:
    """Split text into overlapping chunks on paragraph/sentence boundaries."""
    text = text.strip()
    if not text:
        return []
    # Prefer paragraph boundaries; fall back to a sliding window for long blocks.
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks: list[str] = []
    buf = ""
    for para in paragraphs:
        if len(buf) + len(para) + 1 <= chunk_size:
            buf = f"{buf}\n{para}".strip()
            continue
        if buf:
            chunks.append(buf)
        if len(para) <= chunk_size:
            buf = para
        else:
            for i in range(0, len(para), chunk_size - overlap):
                chunks.append(para[i : i + chunk_size])
            buf = ""
    if buf:
        chunks.append(buf)
    return chunks


class ContextIndexer:
    """Writes context chunks to Elasticsearch, with an in-memory fallback for tests."""

    def __init__(self) -> None:
        self._client = self._init_client()
        self._mem: list[dict] = []

    @property
    def is_memory(self) -> bool:
        return self._client is None

    @staticmethod
    def _init_client():  # type: ignore[no-untyped-def]
        if not settings.elasticsearch_url:
            return None
        try:
            from elasticsearch import Elasticsearch

            kwargs: dict = {"hosts": [settings.elasticsearch_url]}
            if settings.elasticsearch_api_key:
                kwargs["api_key"] = settings.elasticsearch_api_key
            client = Elasticsearch(**kwargs)
            if not client.indices.exists(index=INDEX):
                client.indices.create(
                    index=INDEX,
                    mappings={
                        "properties": {
                            "text": {"type": "text"},
                            "source": {"type": "keyword"},
                            "kind": {"type": "keyword"},
                            "session_id": {"type": "keyword"},
                            "embedding": {
                                "type": "dense_vector",
                                "dims": EMBED_DIM,
                                "index": True,
                                "similarity": "cosine",
                            },
                        }
                    },
                )
            else:
                # 既存 index に session_id keyword mapping を明示する（PR 以前作成の index で
                # session スコープの term フィルタが効くように。冪等 / Codex P2）。
                try:
                    client.indices.put_mapping(
                        index=INDEX, properties={"session_id": {"type": "keyword"}}
                    )
                except Exception as exc:
                    log.warning("ensure_session_id_mapping_failed", error=str(exc))
            return client
        except Exception as exc:  # pragma: no cover - depends on env
            log.warning("elasticsearch_unavailable_using_memory", error=str(exc))
            return None

    def index_context(self, session_id: str, chunks: list[str], source_name: str) -> int:
        """Index `chunks` for a session; returns the number indexed."""
        for i, chunk in enumerate(chunks):
            source = f"{source_name}#{i}"
            text = mask_pii(chunk) if settings.mask_pii_before_index else chunk
            embedding = _embed(text)
            if self._client is not None:  # pragma: no cover - needs live ES
                doc: dict[str, object] = {
                    "text": text,
                    "source": source,
                    "kind": "context",
                    "session_id": session_id,
                }
                if embedding is not None:
                    doc["embedding"] = embedding
                self._client.index(index=INDEX, document=doc)
            else:
                self._mem.append(
                    {"text": text, "source": source, "kind": "context", "session_id": session_id}
                )
        log.info("context_indexed", session=session_id, source=source_name, chunks=len(chunks))
        return len(chunks)

    def delete_repo_context(self, session_id: str) -> int:
        """セッションの GitHub repo 由来 chunk（source が `github:` 始まり）を全削除する。

        準備画面で repo を選び直した / 別 branch へ変えた / 再同期したとき、古い repo・commit の
        コード断片が search_grounding に残って混ざるのを防ぐ（ADR-0028・Codex P2）。repo/path に
        依らず一括で消すため、`delete_context` の `#` 境界一致ではなく `github:` 前方一致で消す。
        削除件数を返す（冪等: 0 件でも安全）。
        """
        if self._client is not None:  # pragma: no cover - needs live ES
            try:
                res = self._client.delete_by_query(
                    index=INDEX,
                    query={
                        "bool": {
                            "filter": [
                                {"term": {"session_id": session_id}},
                                {"prefix": {"source": "github:"}},
                            ]
                        }
                    },
                    refresh=True,
                )
                deleted = int(res.get("deleted", 0))
            except Exception as exc:  # pragma: no cover - depends on env
                log.warning("repo_context_delete_failed", error=str(exc), session=session_id)
                return 0
        else:
            before = len(self._mem)
            self._mem = [
                d
                for d in self._mem
                if not (
                    d.get("session_id") == session_id
                    and str(d.get("source", "")).startswith("github:")
                )
            ]
            deleted = before - len(self._mem)
        if deleted:
            log.info("repo_context_deleted", session=session_id, chunks=deleted)
        return deleted

    def delete_context(self, session_id: str, source_prefix: str) -> int:
        """出所が `source_prefix`（例 `asset:{asset_id}`）の grounding chunk を取り消す（#245）。

        index_context は出所を `{source_name}#{i}` で保存するため、`source_prefix` 自身と
        `source_prefix#*` を削除対象にする（別 source_name への巻き込みを避ける）。ES 接続時は
        delete_by_query、未接続（テスト/ローカル）は in-memory を filter する。削除件数を返す
        （冪等: 0 件でも安全）。これで中断素材の観察を以後の検索（search_grounding）から外す。
        """
        if self._client is not None:  # pragma: no cover - needs live ES
            try:
                # in-memory の _source_matches と同じ `#` 境界にそろえる: 素の prefix（あれば）と
                # `prefix#*` のみを対象にし、別 asset への前方一致の誤爆を防ぐ。
                res = self._client.delete_by_query(
                    index=INDEX,
                    query={
                        "bool": {
                            "filter": [{"term": {"session_id": session_id}}],
                            "minimum_should_match": 1,
                            "should": [
                                {"term": {"source": source_prefix}},
                                {"prefix": {"source": f"{source_prefix}#"}},
                            ],
                        }
                    },
                    refresh=True,
                )
                deleted = int(res.get("deleted", 0))
            except Exception as exc:  # pragma: no cover - depends on env
                log.warning("context_delete_failed", error=str(exc), source=source_prefix)
                return 0
        else:
            before = len(self._mem)
            self._mem = [
                d
                for d in self._mem
                if not (
                    d.get("session_id") == session_id
                    and _source_matches(str(d.get("source", "")), source_prefix)
                )
            ]
            deleted = before - len(self._mem)
        if deleted:
            log.info("context_deleted", session=session_id, source=source_prefix, chunks=deleted)
        return deleted


def _embed(text: str) -> list[float] | None:
    if not (settings.google_api_key or settings.google_genai_use_vertexai):
        return None
    try:  # pragma: no cover - needs creds/network
        from google import genai

        client = genai.Client(api_key=settings.google_api_key or None)
        resp = client.models.embed_content(model=settings.gemini_embed_model, contents=text)
        embeddings = resp.embeddings
        if not embeddings or embeddings[0].values is None:
            return None
        return list(embeddings[0].values)
    except Exception as exc:  # pragma: no cover
        log.warning("embed_failed", error=str(exc))
        return None


class DocumentExtractionError(Exception):
    """バイナリ文書からのテキスト抽出失敗（壊れた zip・想定外の中身・展開爆発）。

    呼び出し側（context/file エンドポイント）はこれを 500 にせず「抽出 0 件」へ平すが、
    メトリクス上は成功（indexed）と区別して計上する（運用での異常検知のため）。
    """


# zip コンテナ（docx/xlsx/pptx）の展開後合計サイズ上限。受理サイズ（max_asset_bytes=25MB）は
# 圧縮後のバイト数なので、極端な圧縮率の zip bomb は展開時にメモリを食い尽くし得る。
# 展開前に infolist の非圧縮サイズ合計で弾く（Cloud Run 同居リクエストを OOM で巻き込まない）。
_MAX_ZIP_EXPANSION_BYTES = 100_000_000


def _guard_zip_expansion(raw: bytes) -> None:
    with zipfile.ZipFile(io.BytesIO(raw)) as archive:
        total = sum(info.file_size for info in archive.infolist())
    if total > _MAX_ZIP_EXPANSION_BYTES:
        raise DocumentExtractionError(
            f"zip expands to {total} bytes (limit {_MAX_ZIP_EXPANSION_BYTES})"
        )


class _HTMLTextExtractor(HTMLParser):
    """HTML から可視テキストだけを取り出す（script/style 等は捨てる）。

    grounding に流すのは「人が読む本文」であって、マークアップや JS ではない。
    依存を増やさず stdlib の HTMLParser で足りる範囲に留める（壊れた HTML でも
    エラーにせず、読めた分だけ返す best-effort）。
    """

    _SKIP_TAGS = frozenset({"script", "style", "noscript", "template"})

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0 and data.strip():
            self._parts.append(data.strip())

    def text(self) -> str:
        return "\n".join(self._parts)


def _extract_pdf(raw: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(raw: bytes) -> str:
    """Word（.docx）の段落と表をテキスト化する。"""
    from docx import Document

    _guard_zip_expansion(raw)
    document = Document(io.BytesIO(raw))
    parts: list[str] = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_xlsx(raw: bytes) -> str:
    """Excel（.xlsx）を全シート TSV 風のテキストにする（数式は計算済み値）。"""
    from openpyxl import load_workbook  # type: ignore[import-untyped]

    _guard_zip_expansion(raw)
    workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells).rstrip())
            if rows:
                parts.append(f"# {sheet.title}")
                parts.extend(rows)
        return "\n".join(parts)
    finally:
        workbook.close()


def _extract_pptx(raw: bytes) -> str:
    """PowerPoint（.pptx）のスライド本文とスピーカーノートをテキスト化する。"""
    from pptx import Presentation  # type: ignore[import-untyped]

    _guard_zip_expansion(raw)
    presentation = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for i, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                texts.append(note)
        if texts:
            parts.append(f"# スライド{i}")
            parts.extend(texts)
    return "\n".join(parts)


def _extract_html(raw: bytes) -> str:
    extractor = _HTMLTextExtractor()
    extractor.feed(raw.decode("utf-8", errors="replace"))
    return extractor.text()


# 拡張子 → 抽出関数。storage.py の TEXT_EXT / DOC_BINARY_EXT（受理判定）とペアで保守する。
_EXTRACTORS: dict[str, Callable[[bytes], str]] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".html": _extract_html,
    ".htm": _extract_html,
}

# MIME → 抽出関数。拡張子なしでも DOC_BINARY_MIME で受理されたファイルを正しく抽出するための
# フォールバック。storage.py の DOC_BINARY_MIME とペアで保守する。
_MIME_EXTRACTORS: dict[str, Callable[[bytes], str]] = {
    "application/pdf": _extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
    "text/html": _extract_html,
}


def extract_text_from_upload(
    filename: str, raw: bytes, content_type: str | None = None
) -> str:
    """Best-effort text extraction (txt/md/html/csv/json/pdf/docx/xlsx/pptx uploads)."""
    name = filename.lower()
    ext = name[name.rfind(".") :] if "." in name else ""
    extractor = _EXTRACTORS.get(ext)
    # 拡張子なしで MIME だけで受理されたバイナリ文書（例: ファイル名 "book"、MIME が docx）は
    # 拡張子ルックアップが失敗するため、MIME でフォールバックする。
    # テキスト系（txt/csv/json）は extractor=None のままで正しく UTF-8 デコードされる。
    if extractor is None and content_type:
        ct = content_type.split(";")[0].strip().lower()
        extractor = _MIME_EXTRACTORS.get(ct)
    if extractor is not None:
        try:
            return extractor(raw)
        except Exception as exc:
            # 壊れたファイル・想定外の中身・zip bomb。呼び出し側が 500 にせず「抽出 0 件」に
            # 平しつつ、成功（indexed）と区別してメトリクス計上できるよう型付き例外で伝える。
            log.warning("document_extract_failed", ext=ext, error=str(exc))
            raise DocumentExtractionError(f"failed to extract {ext}") from exc
    # txt / md / csv / json / anything decodable as utf-8
    return raw.decode("utf-8", errors="replace")
